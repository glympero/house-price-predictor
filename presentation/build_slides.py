"""Build the interview deck from executed model evidence.

Run with ``uv run --group slides python presentation/build_slides.py`` after
training and evaluation. Metrics are read from the persisted JSON artifacts.
"""

import json
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from house_prices import config

OUTPUT = Path(__file__).parent / "slides.pptx"
SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = RGBColor(0x0B, 0x17, 0x2A)
NAVY_2 = RGBColor(0x13, 0x25, 0x3F)
INK = RGBColor(0x16, 0x23, 0x33)
MUTED = RGBColor(0x62, 0x70, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF5, 0xF7, 0xFA)
LINE = RGBColor(0xD9, 0xE0, 0xE8)
BLUE = RGBColor(0x2F, 0x65, 0xCB)
BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFF)
TEAL = RGBColor(0x10, 0xA3, 0x8A)
TEAL_LIGHT = RGBColor(0xE4, 0xF6, 0xF1)
ORANGE = RGBColor(0xF0, 0x9A, 0x2B)
ORANGE_LIGHT = RGBColor(0xFE, 0xF1, 0xDC)
RED = RGBColor(0xD9, 0x4A, 0x4A)
RED_LIGHT = RGBColor(0xFC, 0xE8, 0xE8)
GREEN = RGBColor(0x2D, 0x8A, 0x5E)


def load_results() -> tuple[dict, dict, dict]:
    metadata_path = config.MODELS_DIR / config.METADATA_FILENAME
    evaluation_path = config.MODELS_DIR / config.EVALUATION_FILENAME
    diagnostics_path = config.REPO_ROOT / "docs" / "post_selection_diagnostics.json"
    return (
        json.loads(metadata_path.read_text(encoding="utf-8")),
        json.loads(evaluation_path.read_text(encoding="utf-8")),
        json.loads(diagnostics_path.read_text(encoding="utf-8")),
    )


def blank_slide(prs: Presentation, background: RGBColor = PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = background
    return slide


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.space_after = 0
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor | None = LINE,
    radius: bool = True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_card(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    eyebrow: str,
    value: str,
    detail: str,
    *,
    accent: RGBColor = BLUE,
    fill: RGBColor = WHITE,
):
    add_rect(slide, x, y, width, height, fill=fill)
    add_rect(slide, x, y, 0.08, height, fill=accent, line=None, radius=False)
    compact = height <= 1.05
    add_text(
        slide,
        eyebrow.upper(),
        x + 0.28,
        y + (0.13 if compact else 0.18),
        width - 0.45,
        0.2,
        size=8.5 if compact else 9,
        color=MUTED,
        bold=True,
    )
    add_text(
        slide,
        value,
        x + 0.28,
        y + (0.38 if compact else 0.48),
        width - 0.45,
        0.4,
        size=23 if compact else 25,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        detail,
        x + 0.28,
        y + height - (0.2 if compact else 0.38),
        width - 0.45,
        0.18 if compact else 0.25,
        size=8.5 if compact else 10.5,
        color=MUTED,
    )


def add_header(slide, title: str, kicker: str, number: int, *, dark: bool = False):
    title_color = WHITE if dark else INK
    muted_color = RGBColor(0xA8, 0xB6, 0xC8) if dark else MUTED
    add_text(slide, kicker.upper(), 0.65, 0.38, 8.8, 0.24, size=9.5, color=TEAL, bold=True)
    add_text(slide, title, 0.65, 0.68, 11.5, 0.55, size=27, color=title_color, bold=True)
    add_text(
        slide,
        f"{number:02d}",
        12.15,
        0.43,
        0.5,
        0.3,
        size=10,
        color=muted_color,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide, text: str, *, dark: bool = False):
    color = RGBColor(0x8E, 0x9C, 0xAE) if dark else MUTED
    add_text(slide, text, 0.65, 7.16, 12.0, 0.2, size=8.5, color=color)


def add_label(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    *,
    fill: RGBColor = BLUE_LIGHT,
    color: RGBColor = BLUE,
):
    add_rect(slide, x, y, width, 0.32, fill=fill, line=None)
    add_text(
        slide,
        text.upper(),
        x + 0.08,
        y + 0.08,
        width - 0.16,
        0.15,
        size=8,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_picture_contain(slide, path: Path, x: float, y: float, width: float, height: float):
    with PILImage.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = width / height
    if image_ratio > box_ratio:
        draw_width = width
        draw_height = width / image_ratio
        draw_x = x
        draw_y = y + (height - draw_height) / 2
    else:
        draw_height = height
        draw_width = height * image_ratio
        draw_x = x + (width - draw_width) / 2
        draw_y = y
    return slide.shapes.add_picture(
        str(path),
        Inches(draw_x),
        Inches(draw_y),
        width=Inches(draw_width),
        height=Inches(draw_height),
    )


def add_chevron(slide, x: float, y: float, *, color: RGBColor = LINE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.35), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def model_row(metadata: dict, name: str) -> dict:
    return next(row for row in metadata["cv_comparison"] if row["model"] == name)


def diagnostic_row(diagnostics: dict, name: str) -> dict:
    return next(row for row in diagnostics["model_comparison"] if row["candidate"] == name)


def slide_1(prs: Presentation, metadata: dict, evaluation: dict):
    slide = blank_slide(prs, NAVY)
    add_rect(slide, 0, 0, 0.16, SLIDE_H, fill=TEAL, line=None, radius=False)
    add_text(
        slide,
        "AI SOLUTION ENGINEER · USE CASE 1",
        0.75,
        0.55,
        6.4,
        0.3,
        size=10,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        "House price estimation\nbuilt around evidence",
        0.75,
        1.05,
        7.0,
        1.55,
        size=34,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        (
            "A reproducible supervised-regression workflow: source control, "
            "location-disjoint validation, bounded tuning, protected evaluation "
            "and a deployable analyst API."
        ),
        0.78,
        2.86,
        6.0,
        1.0,
        size=15,
        color=RGBColor(0xBF, 0xCA, 0xD8),
    )

    add_rect(slide, 7.45, 0.75, 5.15, 2.0, fill=NAVY_2, line=None)
    add_text(slide, "SELECTED", 7.82, 1.05, 2.0, 0.25, size=9, color=TEAL, bold=True)
    add_text(
        slide, "Histogram gradient boosting", 7.82, 1.43, 4.25, 0.6, size=24, color=WHITE, bold=True
    )
    add_text(
        slide,
        "Selected under the predefined grouped-CV rule",
        7.82,
        2.2,
        4.1,
        0.28,
        size=11,
        color=RGBColor(0xB7, 0xC4, 0xD4),
    )

    values = [
        ("GROUPED CV", f"{metadata['cv_rmse']:.2f}", "RMSE · original selection", BLUE),
        ("PROTECTED TEST", f"{evaluation['rmse']:.2f}", "RMSE · 83 unseen rows", ORANGE),
        ("LOCATION OVERLAP", str(metadata["n_location_overlap"]), "exact coordinates", TEAL),
    ]
    for index, (label, value, detail, accent) in enumerate(values):
        x = 0.75 + index * 4.05
        add_rect(slide, x, 4.65, 3.65, 1.55, fill=NAVY_2, line=None)
        add_rect(slide, x, 4.65, 3.65, 0.06, fill=accent, line=None, radius=False)
        add_text(
            slide,
            label,
            x + 0.25,
            4.92,
            2.9,
            0.2,
            size=8.5,
            color=RGBColor(0xA8, 0xB6, 0xC8),
            bold=True,
        )
        add_text(slide, value, x + 0.25, 5.18, 1.6, 0.45, size=27, color=WHITE, bold=True)
        add_text(
            slide, detail, x + 1.52, 5.35, 1.8, 0.25, size=9.5, color=RGBColor(0xA8, 0xB6, 0xC8)
        )

    add_text(
        slide,
        "UCI source · SHA-256 pinned · Python 3.12 · FastAPI · Docker",
        0.75,
        6.72,
        8.5,
        0.25,
        size=9,
        color=RGBColor(0x8E, 0x9C, 0xAE),
    )
    add_text(
        slide,
        "01",
        12.1,
        6.72,
        0.5,
        0.25,
        size=9,
        color=RGBColor(0x8E, 0x9C, 0xAE),
        align=PP_ALIGN.RIGHT,
    )


def slide_2(prs: Presentation, metadata: dict):
    slide = blank_slide(prs)
    add_header(slide, "Start with the decision, not the algorithm", "Problem framing", 2)

    steps = [
        ("1", "Evidence", "Historical transactions include inputs and a continuous target."),
        ("2", "Learning task", "Supervised, offline regression for a new supported property."),
        ("3", "User outcome", "An analyst gets a repeatable benchmark, range and caveat."),
    ]
    for index, (number, title, detail) in enumerate(steps):
        y = 1.55 + index * 1.48
        add_rect(slide, 0.7, y, 6.0, 1.12, fill=WHITE)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.95), Inches(y + 0.24), Inches(0.62), Inches(0.62)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE if index < 2 else TEAL
        circle.line.fill.background()
        add_text(
            slide,
            number,
            0.95,
            y + 0.36,
            0.62,
            0.2,
            size=13,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, title, 1.82, y + 0.18, 1.7, 0.3, size=15, bold=True)
        add_text(slide, detail, 1.82, y + 0.53, 4.45, 0.4, size=11, color=MUTED)

    add_rect(slide, 7.15, 1.55, 5.5, 4.95, fill=NAVY, line=None)
    add_label(slide, "Source assumption made explicit", 7.55, 1.92, 2.6, fill=NAVY_2, color=TEAL)
    add_text(
        slide,
        str(metadata["n_dataset_rows"]),
        7.55,
        2.45,
        2.1,
        0.65,
        size=38,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide, "transactions", 9.55, 2.72, 1.5, 0.25, size=12, color=RGBColor(0xB7, 0xC4, 0xD4)
    )
    add_text(
        slide, "Sindian District · 2012–2013", 7.55, 3.25, 4.0, 0.3, size=15, color=WHITE, bold=True
    )
    add_text(slide, "Target", 7.55, 3.88, 1.0, 0.2, size=9, color=TEAL, bold=True)
    add_text(slide, "10,000 TWD per ping", 7.55, 4.14, 3.8, 0.38, size=20, color=WHITE, bold=True)
    add_text(
        slide,
        "Price density—not total property price",
        7.55,
        4.58,
        4.3,
        0.28,
        size=11,
        color=RGBColor(0xB7, 0xC4, 0xD4),
    )
    add_text(
        slide,
        (
            "Matching Kaggle copy was a reasonable identification. The implementation "
            "uses original UCI data plus a pinned checksum."
        ),
        7.55,
        5.28,
        4.4,
        0.75,
        size=12,
        color=RGBColor(0xD3, 0xDB, 0xE5),
    )

    add_footer(
        slide, "Source decision: common dataset name resolved to the original UCI publication"
    )


def slide_3(prs: Presentation, metadata: dict):
    slide = blank_slide(prs)
    add_header(slide, "Location-disjoint development and final test", "Validation design", 3)

    add_card(
        slide,
        0.68,
        1.55,
        2.25,
        1.35,
        "Validated source",
        "414 rows",
        "259 exact locations",
        accent=INK,
    )
    add_chevron(slide, 3.08, 1.95, color=MUTED)
    add_card(
        slide,
        3.58,
        1.35,
        3.45,
        1.75,
        "Development side",
        "331 rows · 207 locations",
        "target-aware EDA + model decisions",
        accent=BLUE,
        fill=BLUE_LIGHT,
    )
    add_chevron(slide, 7.18, 1.95, color=MUTED)
    add_card(
        slide,
        7.67,
        1.35,
        3.3,
        1.75,
        "Protected side",
        "83 rows · 52 locations",
        "opened only after the model was frozen",
        accent=RED,
        fill=RED_LIGHT,
    )
    add_card(
        slide,
        11.25,
        1.55,
        1.42,
        1.35,
        "Overlap",
        "0",
        "coordinate pairs",
        accent=TEAL,
        fill=TEAL_LIGHT,
    )

    add_text(
        slide, "INSIDE THE 331 TRAINING ROWS", 0.75, 3.55, 3.2, 0.2, size=9, color=BLUE, bold=True
    )
    for fold in range(5):
        x = 0.75 + fold * 1.0
        add_rect(slide, x, 3.95, 0.75, 0.75, fill=BLUE_LIGHT, line=BLUE)
        add_text(
            slide,
            f"F{fold + 1}",
            x,
            4.19,
            0.75,
            0.22,
            size=12,
            color=BLUE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_chevron(slide, 5.92, 4.05, color=MUTED)
    add_rect(slide, 6.43, 3.73, 2.6, 1.18, fill=WHITE)
    add_text(
        slide,
        "Bounded GridSearchCV",
        6.68,
        3.98,
        2.1,
        0.27,
        size=14,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "same grouped folds · RMSE",
        6.68,
        4.35,
        2.1,
        0.2,
        size=9.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_chevron(slide, 9.22, 4.05, color=MUTED)
    add_rect(slide, 9.73, 3.73, 2.75, 1.18, fill=TEAL_LIGHT, line=TEAL)
    add_text(
        slide,
        "Freeze lowest CV RMSE",
        9.98,
        3.98,
        2.25,
        0.27,
        size=14,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "then open the protected test",
        9.98,
        4.35,
        2.25,
        0.2,
        size=9.5,
        color=GREEN,
        align=PP_ALIGN.CENTER,
    )

    add_rect(slide, 0.75, 5.45, 11.73, 0.85, fill=NAVY, line=None)
    add_text(
        slide, "Why group coordinates?", 1.05, 5.72, 2.15, 0.25, size=13, color=WHITE, bold=True
    )
    add_text(
        slide,
        (
            "Repeated transactions at one location cannot appear on both sides and "
            "make generalization look easier."
        ),
        3.3,
        5.69,
        8.55,
        0.34,
        size=12.5,
        color=RGBColor(0xD3, 0xDB, 0xE5),
    )

    add_footer(
        slide,
        (
            f"Seed {metadata['random_seed']} · GroupShuffleSplit candidates are "
            f"label-blind · shuffled {metadata['cv_folds']}-fold GroupKFold"
        ),
    )


def slide_4(prs: Presentation):
    slide = blank_slide(prs)
    add_header(
        slide, "Correlation suggests a hypothesis; validation decides", "Training-only EDA", 4
    )

    add_rect(slide, 0.65, 1.4, 7.45, 4.95, fill=WHITE)
    add_picture_contain(
        slide, config.REPORTS_DIR / "eda_mrt_log_transform.png", 0.88, 1.68, 7.0, 4.35
    )
    add_label(slide, "331 training rows only", 1.0, 1.58, 1.65, fill=TEAL_LIGHT, color=TEAL)

    facts = [
        ("−0.67 → −0.71", "MRT correlation after log transform", BLUE),
        ("+0.58", "stores vs price", TEAL),
        ("−0.84", "MRT distance vs longitude", ORANGE),
    ]
    for index, (value, detail, accent) in enumerate(facts):
        y = 1.5 + index * 1.25
        add_card(slide, 8.48, y, 4.15, 1.0, "Pearson association", value, detail, accent=accent)

    add_rect(slide, 8.48, 5.32, 4.15, 1.02, fill=NAVY, line=None)
    add_text(
        slide,
        "association  →  feature hypothesis  →  grouped CV",
        8.72,
        5.57,
        3.7,
        0.23,
        size=11.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Not causal. Correlated predictors can share the same signal.",
        8.72,
        5.92,
        3.7,
        0.2,
        size=9.5,
        color=RGBColor(0xB7, 0xC4, 0xD4),
        align=PP_ALIGN.CENTER,
    )

    add_footer(
        slide,
        "Log-distance improves the linear benchmark; EDA uses development rows only",
    )


def slide_5(prs: Presentation, metadata: dict, diagnostics: dict):
    slide = blank_slide(prs)
    add_header(
        slide,
        "Original selection favored boosting; nested CV narrows the claim",
        "Model evidence",
        5,
    )

    reported_forest = model_row(metadata, "random_forest")
    reported_boosting = model_row(metadata, "gradient_boosting")
    nested_forest = diagnostic_row(diagnostics, "random_forest")
    nested_boosting = diagnostic_row(diagnostics, "gradient_boosting")
    nonlinear_pair = diagnostics["nonlinear_pair"]

    names = [
        ("mean_baseline", "Mean"),
        ("linear_regression", "Raw linear"),
        ("linear_engineered", "Engineered"),
        ("ridge_engineered", "Ridge"),
        ("random_forest", "Forest"),
        ("gradient_boosting", "Boosting"),
    ]
    rows = [model_row(metadata, key) for key, _ in names]
    chart_data = ChartData()
    chart_data.categories = [label for _, label in names]
    chart_data.add_series("Train RMSE", [row["cv_train_rmse"] for row in rows])
    chart_data.add_series("Reported grouped-CV RMSE", [row["cv_rmse"] for row in rows])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.65),
        Inches(1.55),
        Inches(8.1),
        Inches(4.95),
        chart_data,
    ).chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 14
    chart.value_axis.major_unit = 2
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(0xAE, 0xBA, 0xC8)
    chart.series[0].format.line.fill.background()
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = BLUE
    chart.series[1].format.line.fill.background()
    winner_point = chart.series[1].points[len(rows) - 1]
    winner_point.format.fill.solid()
    winner_point.format.fill.fore_color.rgb = TEAL
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0.0"
    plot.data_labels.font.size = Pt(8.5)

    add_rect(slide, 9.05, 1.55, 3.62, 3.15, fill=NAVY, line=None)
    add_label(slide, "Original rule", 9.42, 1.9, 1.55, fill=NAVY_2, color=TEAL)
    add_text(
        slide,
        f"{reported_boosting['cv_rmse']:.2f}",
        9.42,
        2.38,
        1.8,
        0.65,
        size=37,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "boosting reported CV",
        10.58,
        2.68,
        1.55,
        0.23,
        size=9.5,
        color=RGBColor(0xB7, 0xC4, 0xD4),
    )
    add_text(
        slide,
        (f"Boosting {reported_boosting['cv_rmse']:.2f} · Forest {reported_forest['cv_rmse']:.2f}"),
        9.42,
        3.32,
        2.8,
        0.3,
        size=13,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "36-cell vs 18-cell grid · frozen before test",
        9.42,
        3.82,
        2.8,
        0.4,
        size=9.5,
        color=RGBColor(0xB7, 0xC4, 0xD4),
    )

    add_rect(slide, 9.05, 4.94, 3.62, 1.56, fill=ORANGE_LIGHT, line=ORANGE)
    add_text(
        slide,
        "Post-selection nested CV",
        9.35,
        5.16,
        2.9,
        0.25,
        size=11.5,
        color=INK,
        bold=True,
    )
    add_text(
        slide,
        (
            f"Forest {nested_forest['nested_cv_rmse']:.2f} · "
            f"Boosting {nested_boosting['nested_cv_rmse']:.2f}"
        ),
        9.35,
        5.52,
        2.95,
        0.28,
        size=14.5,
        color=ORANGE,
        bold=True,
    )
    add_text(
        slide,
        (
            f"{abs(nonlinear_pair['random_forest_minus_histogram_boosting_mean_rmse']):.2f} "
            "gap · 3–2 folds · no clear winner"
        ),
        9.35,
        5.98,
        2.8,
        0.3,
        size=9.5,
        color=MUTED,
    )

    add_footer(
        slide,
        (
            "The holdout stayed out of nested CV; this diagnostic discloses uncertainty "
            "without reopening selection"
        ),
    )


def slide_6(prs: Presentation):
    slide = blank_slide(prs)
    add_header(
        slide,
        "Linear and tree models use different preprocessing",
        "Model-specific preprocessing",
        6,
    )

    add_text(slide, "INTERPRETABLE BENCHMARK", 0.75, 1.5, 2.4, 0.2, size=9, color=BLUE, bold=True)
    linear_cards = [
        (0.75, "Raw inputs"),
        (2.35, "log MRT + age²"),
        (4.15, "StandardScaler"),
        (5.88, "Ridge · 7.64"),
    ]
    for index, (x, text_value) in enumerate(linear_cards):
        add_rect(slide, x, 1.87, 1.42 if index != 1 else 1.62, 0.72, fill=BLUE_LIGHT, line=BLUE)
        add_text(
            slide,
            text_value,
            x + 0.08,
            2.09,
            (1.26 if index != 1 else 1.46),
            0.25,
            size=10.5,
            color=BLUE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if index < len(linear_cards) - 1:
            add_chevron(slide, x + (1.47 if index != 1 else 1.67), 1.95, color=BLUE)

    add_text(slide, "PRODUCTION PIPELINE", 0.75, 3.2, 2.4, 0.2, size=9, color=TEAL, bold=True)
    production_cards = [
        (0.75, 1.42, "5 raw inputs"),
        (2.55, 1.65, "No scaler\nno redundant transforms"),
        (4.62, 2.45, "Histogram gradient boosting · 6.94"),
    ]
    for index, (x, width, text_value) in enumerate(production_cards):
        add_rect(slide, x, 3.57, width, 0.86, fill=TEAL_LIGHT, line=TEAL)
        add_text(
            slide,
            text_value,
            x + 0.08,
            3.8,
            width - 0.16,
            0.42,
            size=10.5,
            color=GREEN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if index < len(production_cards) - 1:
            add_chevron(slide, x + width + 0.16, 3.72, color=TEAL)

    add_rect(slide, 0.75, 4.92, 7.0, 1.12, fill=NAVY, line=None)
    add_text(slide, "Why different paths?", 1.02, 5.19, 1.75, 0.25, size=12, color=WHITE, bold=True)
    add_text(
        slide,
        (
            "Scaling stabilizes regularized linear optimization. Trees use ordered "
            "thresholds, so monotonic log/square transforms add no split information."
        ),
        2.68,
        5.12,
        4.66,
        0.55,
        size=11.5,
        color=RGBColor(0xD3, 0xDB, 0xE5),
    )

    add_rect(slide, 8.15, 1.5, 4.48, 4.54, fill=WHITE)
    add_picture_contain(
        slide, config.REPORTS_DIR / "gradient_descent_reference.png", 8.35, 1.72, 4.08, 2.85
    )
    add_label(slide, "Implementation check", 8.48, 1.72, 1.85, fill=ORANGE_LIGHT, color=ORANGE)
    add_text(
        slide,
        "Manual GD and scikit-learn converge to the same fit",
        8.48,
        4.72,
        3.75,
        0.42,
        size=12.5,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "manual   w = -9.2234  ·  b = 37.6697",
        8.48,
        5.18,
        3.75,
        0.22,
        size=10.5,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "sklearn  w = -9.2234  ·  b = 37.6698",
        8.48,
        5.48,
        3.75,
        0.22,
        size=10.5,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "331 training rows · standardized log-MRT · agreement within 0.0001",
        8.48,
        5.78,
        3.75,
        0.2,
        size=8.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    add_footer(
        slide,
        (
            "One-feature linear verification only; the frozen histogram-boosting model "
            "was trained separately"
        ),
    )


def slide_7(prs: Presentation, evaluation: dict, diagnostics: dict):
    slide = blank_slide(prs)
    add_header(slide, "Final performance on unseen locations", "Protected holdout", 7)

    sensitivity = diagnostics["protected_holdout_sensitivity"]

    add_label(
        slide,
        "83 rows · 52 unseen locations · zero overlap",
        0.72,
        1.42,
        3.75,
        fill=RED_LIGHT,
        color=RED,
    )
    metrics = [
        (
            "RMSE",
            f"{evaluation['rmse']:.2f}",
            f"≈ {round(evaluation['rmse'] * 10_000, -2):,.0f} TWD/ping",
            ORANGE,
        ),
        (
            "MAE",
            f"{evaluation['mae']:.2f}",
            f"≈ {round(evaluation['mae'] * 10_000, -2):,.0f} TWD/ping avg. miss",
            BLUE,
        ),
        ("R²", f"{evaluation['r2']:.3f}", "56.9% less squared error vs mean", TEAL),
    ]
    for index, (label, value, unit, accent) in enumerate(metrics):
        add_card(
            slide,
            0.72,
            1.96 + index * 1.16,
            3.58,
            0.92,
            label,
            value,
            unit,
            accent=accent,
        )

    ci = evaluation["rmse_bootstrap_95_ci"]
    add_rect(slide, 0.72, 5.55, 3.58, 0.86, fill=NAVY, line=None)
    add_text(
        slide,
        "RMSE bootstrap 95% interval",
        0.98,
        5.75,
        2.95,
        0.2,
        size=9.5,
        color=RGBColor(0xB7, 0xC4, 0xD4),
        bold=True,
    )
    add_text(
        slide, f"{ci[0]:.2f} — {ci[1]:.2f}", 0.98, 6.05, 2.95, 0.25, size=16, color=WHITE, bold=True
    )

    add_rect(slide, 4.62, 1.42, 8.0, 4.8, fill=WHITE)
    add_picture_contain(
        slide, config.REPORTS_DIR / "model_pred_vs_actual.png", 4.88, 1.68, 7.5, 4.25
    )
    add_rect(slide, 7.12, 5.52, 5.12, 0.94, fill=RED_LIGHT, line=RED)
    add_text(
        slide,
        (
            f"117.5 actual → 40.1 predicted · "
            f"{sensitivity['largest_error_share_of_total_squared_error']:.0%} of squared error"
        ),
        7.34,
        5.72,
        4.7,
        0.23,
        size=10.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        (
            f"Training max {sensitivity['development_target_max']:.1f} · "
            f"RMSE without that row {sensitivity['rmse_without_largest_error']:.2f}"
        ),
        7.34,
        6.08,
        4.7,
        0.2,
        size=9,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_footer(
        slide,
        (
            "This explains the estimate's imprecision; it does not excuse the miss or "
            "replace the official 10.48"
        ),
    )


def slide_8(prs: Presentation, metadata: dict, evaluation: dict):
    slide = blank_slide(prs)
    add_header(
        slide,
        "Error patterns and uncertainty",
        "Diagnostics and uncertainty",
        8,
    )

    add_rect(slide, 0.65, 1.45, 6.0, 3.78, fill=WHITE)
    add_picture_contain(
        slide, config.REPORTS_DIR / "validation_error_by_target_band.png", 0.92, 1.75, 5.45, 2.85
    )
    add_text(
        slide,
        "Training-only grouped OOF error",
        0.95,
        4.64,
        5.35,
        0.25,
        size=11.5,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "High-price RMSE 9.32 vs 5.58 / 5.40",
        0.95,
        4.94,
        5.35,
        0.2,
        size=9.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_rect(slide, 6.93, 1.45, 5.72, 3.78, fill=WHITE)
    add_picture_contain(
        slide, config.REPORTS_DIR / "model_permutation_importance.png", 7.2, 1.75, 5.18, 2.85
    )
    add_text(
        slide,
        "Protected-holdout permutation importance",
        7.22,
        4.64,
        5.12,
        0.25,
        size=11.5,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Descriptive; correlated inputs can share importance",
        7.22,
        4.94,
        5.12,
        0.2,
        size=9.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    bounds = metadata["residual_bounds"]
    add_rect(slide, 0.65, 5.55, 12.0, 0.92, fill=NAVY, line=None)
    add_text(
        slide,
        "90% residual range",
        0.95,
        5.83,
        1.6,
        0.22,
        size=10,
        color=RGBColor(0xB7, 0xC4, 0xD4),
        bold=True,
    )
    add_text(
        slide,
        f"{bounds['lower']:+.2f}",
        2.58,
        5.78,
        0.85,
        0.3,
        size=16,
        color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_rect(slide, 3.47, 5.91, 3.4, 0.08, fill=TEAL, line=None, radius=False)
    add_text(
        slide,
        "prediction",
        4.52,
        5.72,
        1.3,
        0.2,
        size=9,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        f"{bounds['upper']:+.2f}",
        6.9,
        5.78,
        0.85,
        0.3,
        size=16,
        color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        f"Observed coverage {evaluation['interval_coverage']:.1%}",
        8.15,
        5.78,
        2.15,
        0.3,
        size=15,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "finite sample · fixed width · not an individual guarantee",
        10.05,
        5.83,
        2.25,
        0.22,
        size=9,
        color=RGBColor(0xB7, 0xC4, 0xD4),
        align=PP_ALIGN.RIGHT,
    )

    add_footer(
        slide,
        (
            "Interpretation combines training-only error slices with descriptive "
            "final-test diagnostics"
        ),
    )


def slide_9(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "One frozen pipeline from browser to prediction", "Solution architecture", 9)

    boxes = [
        (0.72, "Analyst UI", "guided scenarios\nraw request/response", BLUE_LIGHT, BLUE),
        (3.08, "FastAPI", "Pydantic schema\nmodel lifecycle", TEAL_LIGHT, TEAL),
        (5.44, "Support gate", "training ranges\n422 if unsupported", ORANGE_LIGHT, ORANGE),
        (7.8, "Frozen pipeline", "raw features\nboosted trees", BLUE_LIGHT, BLUE),
        (10.16, "Response", "point + interval\nunit + caveat", TEAL_LIGHT, TEAL),
    ]
    for index, (x, title, detail, fill, accent) in enumerate(boxes):
        add_rect(slide, x, 1.72, 2.02, 1.55, fill=fill, line=accent)
        add_text(
            slide,
            title,
            x + 0.13,
            2.04,
            1.76,
            0.28,
            size=14,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide, detail, x + 0.13, 2.48, 1.76, 0.46, size=10, color=MUTED, align=PP_ALIGN.CENTER
        )
        if index < len(boxes) - 1:
            add_chevron(slide, x + 2.07, 2.2, color=MUTED)

    add_rect(slide, 0.72, 3.86, 3.7, 1.45, fill=NAVY, line=None)
    add_text(slide, "Lifecycle contract", 1.02, 4.15, 1.8, 0.25, size=13, color=WHITE, bold=True)
    add_text(
        slide,
        "/health = process alive\n/ready = model loaded or 503",
        1.02,
        4.55,
        2.9,
        0.5,
        size=12,
        color=RGBColor(0xD3, 0xDB, 0xE5),
    )

    add_rect(slide, 4.7, 3.86, 3.7, 1.45, fill=WHITE)
    add_text(slide, "Container", 5.0, 4.15, 1.4, 0.25, size=13, bold=True)
    add_text(
        slide,
        "Python 3.12 slim · locked deps\nUID 10001 · model built in image",
        5.0,
        4.55,
        2.9,
        0.5,
        size=11,
        color=MUTED,
    )

    add_rect(slide, 8.68, 3.86, 3.97, 1.45, fill=WHITE)
    add_text(slide, "CI evidence", 8.98, 4.15, 1.5, 0.25, size=13, bold=True)
    add_text(
        slide,
        "Ruff · tests · train/evaluate\nDocker · readiness · real prediction",
        8.98,
        4.55,
        3.05,
        0.5,
        size=11,
        color=MUTED,
    )

    add_rect(slide, 0.72, 5.75, 11.93, 0.64, fill=TEAL_LIGHT, line=TEAL)
    add_text(
        slide,
        (
            "Artifact evidence: model + best parameters + split/fold protocol + "
            "support ranges + residual offsets + code/data hashes"
        ),
        1.02,
        5.96,
        11.3,
        0.22,
        size=11,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_footer(
        slide,
        (
            "Stateless prototype now; registry, feedback store and multi-version "
            "rollout are explicit extension seams"
        ),
    )


def slide_10(prs: Presentation, diagnostics: dict):
    slide = blank_slide(prs, NAVY)
    add_header(slide, "The interview story in five moves", "Close", 10, dark=True)

    sensitivity = diagnostics["protected_holdout_sensitivity"]

    story = [
        ("1", "Frame", "Supervised regression for analyst decision support."),
        ("2", "Protect", "Pin the source; split locations before target-aware EDA."),
        ("3", "Learn", "Use correlations for hypotheses; compare on grouped folds."),
        ("4", "Choose", "Predeclare the rule; freeze the grouped-CV winner."),
        ("5", "Report", "Disclose nested-CV ambiguity and expensive-tail weakness."),
    ]
    for index, (number, title, detail) in enumerate(story):
        y = 1.48 + index * 0.9
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.78), Inches(y), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        add_text(
            slide,
            number,
            0.78,
            y + 0.14,
            0.5,
            0.18,
            size=10,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, title, 1.55, y + 0.03, 1.05, 0.25, size=14, color=WHITE, bold=True)
        add_text(
            slide, detail, 2.65, y + 0.04, 4.3, 0.38, size=11.5, color=RGBColor(0xC7, 0xD1, 0xDE)
        )

    add_rect(slide, 7.45, 1.48, 5.12, 2.05, fill=NAVY_2, line=None)
    add_text(slide, "Current limits", 7.8, 1.82, 2.0, 0.3, size=17, color=WHITE, bold=True)
    add_text(
        slide,
        (
            "• 414 rows · one district · 2012–2013\n"
            "• nonlinear ranking is uncertain under nested CV\n"
            f"• one row drives {sensitivity['largest_error_share_of_total_squared_error']:.0%} "
            "of final squared error"
        ),
        7.8,
        2.3,
        4.15,
        0.9,
        size=12,
        color=RGBColor(0xD3, 0xDB, 0xE5),
    )

    add_rect(slide, 7.45, 3.82, 5.12, 1.72, fill=TEAL_LIGHT, line=None)
    add_text(slide, "Next step", 7.8, 4.13, 2.3, 0.28, size=16, color=GREEN, bold=True)
    add_text(
        slide,
        (
            "Predeclare nested selection on broader, recent data—especially expensive "
            "properties—then evaluate once on a fresh time/geography holdout."
        ),
        7.8,
        4.58,
        4.1,
        0.62,
        size=12,
        color=INK,
    )

    add_rect(slide, 0.75, 6.2, 11.82, 0.62, fill=NAVY_2, line=None)
    add_text(
        slide,
        (
            "AI coding assistants supported scaffolding, tests, documentation and "
            "validation review; results were verified through executed notebooks, "
            "hashes, tests and container checks."
        ),
        1.0,
        6.41,
        11.3,
        0.22,
        size=9.5,
        color=RGBColor(0xA8, 0xB6, 0xC8),
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Questions", dark=True)


def build(metadata: dict, evaluation: dict, diagnostics: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide_1(prs, metadata, evaluation)
    slide_2(prs, metadata)
    slide_3(prs, metadata)
    slide_4(prs)
    slide_5(prs, metadata, diagnostics)
    slide_6(prs)
    slide_7(prs, evaluation, diagnostics)
    slide_8(prs, metadata, evaluation)
    slide_9(prs)
    slide_10(prs, diagnostics)
    return prs


def main() -> None:
    metadata, evaluation, diagnostics = load_results()
    prs = build(metadata, evaluation, diagnostics)
    prs.save(OUTPUT)
    print(f"wrote {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
